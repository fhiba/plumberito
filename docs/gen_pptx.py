from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand tokens ──────────────────────────────────────────────────────────────
CREAM    = RGBColor(0xFF, 0xF8, 0xEF)
DARK     = RGBColor(0x1E, 0x1B, 0x13)
RED      = RGBColor(0xCC, 0x00, 0x00)
BROWN    = RGBColor(0x7A, 0x52, 0x30)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

W = Inches(13.33)   # 16:9 widescreen
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(3)):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=Pt(18), bold=False, color=DARK,
             align=PP_ALIGN.LEFT, wrap=True, font_name="Arial"):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def slide_bg(slide, color=CREAM):
    add_rect(slide, 0, 0, W, H, fill=color)

def label(slide, text, l, t, color=RED):
    add_text(slide, text, l, t, Inches(10), Inches(0.35),
             font_size=Pt(10), bold=True, color=color,
             font_name="Courier New")

# ── SLIDE 1 — Title ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s, DARK)

# Red accent bar left
add_rect(s, 0, 0, Inches(0.18), H, fill=RED, line_w=Pt(0))

# PLUMBERITO wordmark
add_text(s, "PLUMBERITO", Inches(0.5), Inches(1.8), Inches(12), Inches(2.2),
         font_size=Pt(96), bold=True, color=CREAM,
         align=PP_ALIGN.LEFT, font_name="Arial Black")

# Tagline
add_text(s, "Ship. Deploy. Monitor. Fix — con un solo prompt.",
         Inches(0.55), Inches(3.9), Inches(10), Inches(0.7),
         font_size=Pt(22), bold=False, color=RGBColor(0xCC,0xBB,0xAA),
         font_name="Courier New")

# Bottom bar
add_rect(s, 0, H - Inches(0.5), W, Inches(0.5), fill=RED, line_w=Pt(0))
add_text(s, "HackITBA 2026", Inches(0.4), H - Inches(0.42), Inches(6), Inches(0.4),
         font_size=Pt(11), bold=True, color=CREAM, font_name="Courier New")

# ── SLIDE 2 — El Problema ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_rect(s, 0, 0, Inches(0.18), H, fill=RED, line_w=Pt(0))

label(s, "01 / PROBLEMÁTICA", Inches(0.4), Inches(0.3))
add_text(s, "Demasiadas herramientas.\nCero cohesión.", Inches(0.4), Inches(0.75),
         Inches(8), Inches(1.6), font_size=Pt(44), bold=True, color=DARK,
         font_name="Arial Black")

problems = [
    ("🔧", "5+ servicios separados", "GitHub, Vercel, AWS, Sentry, Railway... cada uno con su curva de aprendizaje."),
    ("⏱", "Setup que tarda horas",   "Configurar CI/CD, IAM, DNS y entornos consume más tiempo que escribir código."),
    ("💸", "Costos opacos",           "Cada plataforma cobra por separado. El total explota sin que lo notes."),
    ("🚨", "Incidents sin contexto",  "Un error en prod requiere abrir 4 tabs para entender qué pasó."),
]

for i, (icon, title, desc) in enumerate(problems):
    col = 0 if i < 2 else 1
    row = i % 2
    lx = Inches(0.4) + col * Inches(6.4)
    ty = Inches(2.5) + row * Inches(2.1)
    box = add_rect(s, lx, ty, Inches(6.0), Inches(1.8), fill=WHITE, line=DARK, line_w=Pt(2.5))
    add_text(s, f"{icon}  {title}", lx + Inches(0.15), ty + Inches(0.1),
             Inches(5.7), Inches(0.5), font_size=Pt(15), bold=True, color=RED, font_name="Arial")
    add_text(s, desc, lx + Inches(0.15), ty + Inches(0.55),
             Inches(5.7), Inches(1.1), font_size=Pt(12), color=DARK, font_name="Courier New")

# ── SLIDE 3 — La Solución ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s, DARK)
add_rect(s, 0, 0, Inches(0.18), H, fill=RED, line_w=Pt(0))

label(s, "02 / SOLUCIÓN", Inches(0.4), Inches(0.3), color=RGBColor(0xCC,0xBB,0xAA))
add_text(s, "Un solo lugar.\nUn solo prompt.", Inches(0.4), Inches(0.75),
         Inches(9), Inches(1.8), font_size=Pt(48), bold=True, color=CREAM,
         font_name="Arial Black")

add_text(s,
    "Plumberito es un agente de IA que toma tu idea y, desde una sola interfaz,\n"
    "genera el código, crea el repo, provisiona la infraestructura, deployea la app\n"
    "y monitorea errores — respondiendo con un PR automático si algo falla en prod.",
    Inches(0.4), Inches(2.7), Inches(12.5), Inches(2.0),
    font_size=Pt(16), color=RGBColor(0xCC,0xBB,0xAA), font_name="Courier New")

# Pipeline visual
steps = ["DESCRIBE", "GENERATE", "DEPLOY", "MONITOR", "FIX"]
colors_bg = [DARK, DARK, DARK, DARK, RED]
for i, (step, bg) in enumerate(zip(steps, colors_bg)):
    lx = Inches(0.4) + i * Inches(2.52)
    ty = Inches(5.1)
    bw = Inches(2.3)
    bh = Inches(1.6)
    add_rect(s, lx, ty, bw, bh, fill=bg, line=CREAM, line_w=Pt(2))
    add_text(s, f"{i+1:02d}", lx + Inches(0.12), ty + Inches(0.1),
             bw, Inches(0.4), font_size=Pt(10), bold=True,
             color=RGBColor(0xCC,0xBB,0xAA), font_name="Courier New")
    add_text(s, step, lx + Inches(0.12), ty + Inches(0.5),
             bw, Inches(0.6), font_size=Pt(18), bold=True,
             color=CREAM, font_name="Arial Black")
    # Arrow
    if i < len(steps) - 1:
        add_text(s, "→", lx + bw - Inches(0.05), ty + Inches(0.55),
                 Inches(0.35), Inches(0.5), font_size=Pt(20), color=RED,
                 font_name="Arial", align=PP_ALIGN.CENTER)

# ── SLIDE 4 — Incident Management ────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_rect(s, 0, 0, Inches(0.18), H, fill=RED, line_w=Pt(0))

label(s, "03 / INCIDENT MANAGEMENT", Inches(0.4), Inches(0.3))
add_text(s, "De Sentry al fix,\nen segundos.", Inches(0.4), Inches(0.75),
         Inches(8), Inches(1.8), font_size=Pt(44), bold=True, color=DARK,
         font_name="Arial Black")

flow = [
    ("🔴  Error en prod", "Sentry detecta la excepción y notifica a Plumberito."),
    ("🧠  Análisis con IA", "El agente lee el stack trace, localiza el root cause en el repo."),
    ("🔧  Patch automático", "Genera un fix, lo commitea en una rama nueva."),
    ("📬  PR a tu repo", "Abre el Pull Request con la explicación del bug y el parche listo para review."),
]

for i, (title, desc) in enumerate(flow):
    ty = Inches(2.6) + i * Inches(1.1)
    # Number circle
    add_rect(s, Inches(0.4), ty, Inches(0.6), Inches(0.75), fill=RED, line_w=Pt(0))
    add_text(s, str(i+1), Inches(0.4), ty + Inches(0.1), Inches(0.6), Inches(0.55),
             font_size=Pt(18), bold=True, color=CREAM, align=PP_ALIGN.CENTER, font_name="Arial Black")
    add_text(s, title, Inches(1.15), ty, Inches(5), Inches(0.4),
             font_size=Pt(14), bold=True, color=RED, font_name="Arial")
    add_text(s, desc, Inches(1.15), ty + Inches(0.38), Inches(11.5), Inches(0.6),
             font_size=Pt(12), color=DARK, font_name="Courier New")

# Right side: "resultado"
add_rect(s, Inches(7.3), Inches(2.4), Inches(5.6), Inches(4.6), fill=DARK, line=RED, line_w=Pt(2))
add_text(s, "# resultado", Inches(7.5), Inches(2.6), Inches(5.2), Inches(0.4),
         font_size=Pt(11), bold=True, color=RGBColor(0x7A,0x52,0x30), font_name="Courier New")
code = (
    "PR #42 — auto-fix: null ref\n"
    "en UserService.getProfile()\n\n"
    "- return user.profile\n"
    "+ return user?.profile ?? {}\n\n"
    "✓ Tests pasando\n"
    "✓ Sin breaking changes\n"
    "→ Listo para merge"
)
add_text(s, code, Inches(7.5), Inches(3.1), Inches(5.2), Inches(3.7),
         font_size=Pt(12), color=CREAM, font_name="Courier New")

# ── SLIDE 5 — TAM / SAM / SOM ────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
add_rect(s, 0, 0, Inches(0.18), H, fill=RED, line_w=Pt(0))

label(s, "04 / MERCADO", Inches(0.4), Inches(0.3))
add_text(s, "El mercado es enorme\ny está fragmentado.", Inches(0.4), Inches(0.75),
         Inches(8), Inches(1.4), font_size=Pt(38), bold=True, color=DARK,
         font_name="Arial Black")

# Concentric circles (simulated with nested rectangles for pptx simplicity)
# TAM — big
cx, cy = Inches(4.2), Inches(4.5)
r_tam = Inches(2.9)
r_sam = Inches(2.1)
r_som = Inches(1.2)

from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

def add_circle(slide, cx, cy, r, fill, line, alpha=None):
    shape = slide.shapes.add_shape(9, cx - r, cy - r, r*2, r*2)  # 9 = oval
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(2.5)
    return shape

c_tam = add_circle(s, cx, cy, r_tam, RGBColor(0xF0,0xE8,0xDF), DARK)
c_sam = add_circle(s, cx, cy, r_sam, RGBColor(0xD4,0xBB,0xA0), DARK)
c_som = add_circle(s, cx, cy, r_som, RED, RED)

# Labels inside circles
add_text(s, "SOM\n$320M", cx - Inches(0.8), cy - Inches(0.45),
         Inches(1.6), Inches(0.9), font_size=Pt(13), bold=True,
         color=CREAM, align=PP_ALIGN.CENTER, font_name="Arial Black")

add_text(s, "SAM", cx - Inches(1.6), cy - r_sam + Inches(0.15),
         Inches(3.2), Inches(0.4), font_size=Pt(11), bold=True,
         color=DARK, align=PP_ALIGN.CENTER, font_name="Arial Black")

add_text(s, "TAM", cx - r_tam + Inches(0.1), cy - r_tam + Inches(0.15),
         r_tam * 2 - Inches(0.2), Inches(0.4), font_size=Pt(11), bold=True,
         color=DARK, align=PP_ALIGN.CENTER, font_name="Arial Black")

# Right side — detail boxes
details = [
    ("TAM", "$28B", "Mercado global de\nherramientas de desarrollo"),
    ("SAM", "$4.5B", "DevOps tooling para\nstartups y equipos chicos"),
    ("SOM", "$320M", "Indie devs + freelancers\nLatam & global"),
]
for i, (tag, val, desc) in enumerate(details):
    lx = Inches(8.0)
    ty = Inches(1.8) + i * Inches(1.7)
    add_rect(s, lx, ty, Inches(4.9), Inches(1.5), fill=WHITE, line=DARK, line_w=Pt(2))
    add_text(s, tag, lx + Inches(0.15), ty + Inches(0.1),
             Inches(1.2), Inches(0.45), font_size=Pt(13), bold=True,
             color=RED, font_name="Arial Black")
    add_text(s, val, lx + Inches(1.3), ty + Inches(0.05),
             Inches(3.4), Inches(0.55), font_size=Pt(24), bold=True,
             color=DARK, font_name="Arial Black")
    add_text(s, desc, lx + Inches(0.15), ty + Inches(0.7),
             Inches(4.6), Inches(0.7), font_size=Pt(11), color=DARK,
             font_name="Courier New")

# ── SLIDE 6 — Modelo de Negocio ──────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s, DARK)
add_rect(s, 0, 0, Inches(0.18), H, fill=RED, line_w=Pt(0))

label(s, "05 / MODELO DE NEGOCIO", Inches(0.4), Inches(0.3), color=RGBColor(0xCC,0xBB,0xAA))
add_text(s, "Freemium → Pro", Inches(0.4), Inches(0.75),
         Inches(10), Inches(1.2), font_size=Pt(52), bold=True, color=CREAM,
         font_name="Arial Black")

tiers = [
    {
        "name": "FREE",
        "price": "$0 / mes",
        "color": DARK,
        "border": RGBColor(0xCC,0xBB,0xAA),
        "items": ["3 deployments / mes", "Infraestructura básica", "Monitoreo limitado", "Repositorio público"],
        "highlight": False,
    },
    {
        "name": "PRO",
        "price": "$29 / mes",
        "color": RED,
        "border": RED,
        "items": ["Deployments ilimitados", "Incident management + PR auto-fix", "Monitoreo 24/7 con Sentry", "Repos privados + team collab"],
        "highlight": True,
    },
    {
        "name": "TEAM",
        "price": "$99 / mes",
        "color": DARK,
        "border": BROWN,
        "items": ["Todo lo de Pro", "Hasta 10 usuarios", "SLA de respuesta", "Dashboard de costos unificado"],
        "highlight": False,
    },
]

for i, tier in enumerate(tiers):
    lx = Inches(0.45) + i * Inches(4.27)
    ty = Inches(2.1)
    bw = Inches(4.0)
    bh = Inches(4.9)
    bg = tier["color"] if tier["highlight"] else RGBColor(0x2A, 0x25, 0x1C)
    add_rect(s, lx, ty, bw, bh, fill=bg, line=tier["border"], line_w=Pt(3))
    add_text(s, tier["name"], lx + Inches(0.2), ty + Inches(0.2),
             bw - Inches(0.3), Inches(0.55), font_size=Pt(20), bold=True,
             color=CREAM, font_name="Arial Black")
    add_text(s, tier["price"], lx + Inches(0.2), ty + Inches(0.85),
             bw - Inches(0.3), Inches(0.65), font_size=Pt(26), bold=True,
             color=CREAM if tier["highlight"] else RGBColor(0xCC,0xBB,0xAA),
             font_name="Arial Black")
    for j, item in enumerate(tier["items"]):
        add_text(s, f"→  {item}", lx + Inches(0.2), ty + Inches(1.7) + j * Inches(0.72),
                 bw - Inches(0.3), Inches(0.6), font_size=Pt(11),
                 color=CREAM, font_name="Courier New")

# ── SLIDE 7 — Demo ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s, DARK)
add_rect(s, 0, 0, W, Inches(0.18), fill=RED, line_w=Pt(0))
add_rect(s, 0, H - Inches(0.18), W, Inches(0.18), fill=RED, line_w=Pt(0))

add_text(s, "DEMO", Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.0),
         font_size=Pt(140), bold=True, color=CREAM,
         align=PP_ALIGN.CENTER, font_name="Arial Black")

add_text(s, "A continuación — demo en vivo del producto",
         Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.7),
         font_size=Pt(20), color=RGBColor(0xCC,0xBB,0xAA),
         align=PP_ALIGN.CENTER, font_name="Courier New")

add_text(s, "PLUMBERITO  ·  HackITBA 2026",
         Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
         font_size=Pt(12), bold=True, color=RED,
         align=PP_ALIGN.CENTER, font_name="Courier New")

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/home/jaiba/Documents/plumberito/docs/Plumberito_HackITBA2026.pptx"
prs.save(out)
print(f"Saved → {out}")
